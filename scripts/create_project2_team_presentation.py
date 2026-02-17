from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


TITLE = "Bank Lending Strategy Optimizer"
SUBTITLE = "Project 2 Team Presentation (Requirements-Aligned)"


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)


def add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    left, top, width, height = Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.6)
    tbl = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table

    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            tbl.cell(r, c).text = value


def build_presentation(output_path: Path) -> None:
    prs = Presentation()

    add_title_slide(
        prs,
        TITLE,
        f"{SUBTITLE}\nNode + Express + Handlebars + PostgreSQL",
    )

    add_bullets_slide(prs, "Elevator Pitch", [
        "A full-stack decision-support app for comparing bank performance under changing macro conditions.",
        "Combines FDIC institution data with FRED economic indicators in one workflow.",
        "Adds strategy CRUD and analysis summaries to support lending portfolio decisions.",
    ])

    add_bullets_slide(prs, "Problem and Concept", [
        "Banks need better visibility into performance trends and economic context.",
        "Static spreadsheets do not support fast comparison or scenario thinking.",
        "Our app delivers interactive dashboards + strategy notes + analysis artifacts.",
    ])

    add_bullets_slide(prs, "Architecture and Stack", [
        "Backend: Node.js + Express.js with REST APIs.",
        "Frontend: Handlebars-rendered pages with interactive UI components.",
        "Database: PostgreSQL (SQL-first via pg) with multi-table schema.",
        "Structure: MVC pattern across routes/controllers/views/models.",
        "Config: environment variables via .env and dotenv.",
    ])

    add_bullets_slide(prs, "Core Features Implemented", [
        "Bank Performance dashboard with metrics and comparison views.",
        "US Economic Performance page using engineered macro indicators.",
        "Full Strategy CRUD (GET/POST/PUT/DELETE) for bank notes and strategy records.",
        "About + Methodology & Analysis tab with model and clustering outputs.",
    ])

    add_bullets_slide(prs, "Data and APIs", [
        "FDIC BankFind API for bank financial/institution data.",
        "FRED API for rates, unemployment, GDP, and related macro signals.",
        "Database tables include users, economic_data, bank_performance, and strategy-related tables.",
        "CSV exports support reproducibility and instructor reruns.",
    ])

    add_bullets_slide(prs, "Analysis Highlights", [
        "Best holdout model: winsorized linear regression (RMSE 0.3531).",
        "XGBoost and Prophet showed larger train-vs-test gaps in this setup.",
        "KMeans clustering with PCA projection identifies 3 peer-group segments.",
        "Methodology visualizations are now embedded directly in the web app.",
    ])

    requirement_rows = [
        ["Node + Express", "✅ Met", "Express app with structured API + page routes"],
        ["Handlebars", "✅ Met", "Server-side rendering across core pages"],
        ["PostgreSQL", "✅ Met", "SQL schema + seeded data + live queries"],
        ["MVC structure", "✅ Met", "Routes/controllers/views/models folders in use"],
        ["GET/POST/PUT/DELETE", "✅ Met", "Strategies endpoints implement full CRUD"],
        ["Environment variables", "✅ Met", ".env pattern + dotenv configuration"],
        ["2+ tables, 5+ records", "✅ Met", "Multiple domain tables with seeded records"],
        ["Authentication", "⚠ In Progress", "User table/session scaffolding exists; full auth UI pending"],
        ["Render deployment", "⚠ In Progress", "Local app complete; deployment step queued"],
    ]
    add_table_slide(
        prs,
        "Project 2 Team Requirement Coverage",
        ["Requirement", "Status", "Evidence in This Project"],
        requirement_rows,
    )

    add_bullets_slide(prs, "Exploration Requirement (2+)", [
        "External APIs: FDIC + FRED integration.",
        "New packages/tooling: pg, express-handlebars, express-session, dotenv, bcrypt.",
        "Advanced analytics integration: XGBoost, Prophet, KMeans/PCA methodology content.",
    ])

    add_bullets_slide(prs, "Process and Collaboration", [
        "Iterative build: schema -> seed pipeline -> API routes -> UI pages -> analysis tab.",
        "Debug cycle included route validation, server process cleanup, and UX refinements.",
        "Focus on presentation readiness and requirement traceability.",
    ])

    add_bullets_slide(prs, "Demo Flow", [
        "1) Open Bank Performance and Economic dashboards.",
        "2) Show strategy CRUD workflow in Bank Notes.",
        "3) Open Methodology & Analysis tab for code/graphs and PCA cluster scatter.",
        "4) Close with requirements matrix and roadmap.",
    ])

    add_bullets_slide(prs, "Future Development", [
        "Complete full user authentication flow and role-protected actions.",
        "Deploy on Render and add production configuration checks.",
        "Improve chart aesthetics further and add downloadable report summaries.",
        "Expand model monitoring and scenario simulation endpoints.",
    ])

    add_bullets_slide(prs, "Submission Checklist", [
        "Deployed app URL (add after Render publish)",
        "GitHub repository URL",
        "README with screenshots + deployed link",
        "Team presentation: each member speaks at least once",
    ])

    prs.save(output_path)


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[1] / "Project2_Team_Presentation_Bank_Lending_Strategy_Optimizer.pptx"
    build_presentation(out)
    print(f"Created: {out}")
